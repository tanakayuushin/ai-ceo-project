from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# カラー定義
NAVY = RGBColor(0x1E, 0x3A, 0x5F)
GOLD = RGBColor(0xC9, 0xA8, 0x4C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CHARCOAL = RGBColor(0x33, 0x33, 0x33)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
GREEN = RGBColor(0x27, 0xAE, 0x60)
RED = RGBColor(0xE7, 0x4C, 0x3C)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width = W
prs.slide_height = H

blank_layout = prs.slide_layouts[6]  # blank


def add_rect(slide, left, top, width, height, fill_color=None, line_color=None):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, left, top, width, height,
             font_size=24, bold=False, color=CHARCOAL,
             align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox


def slide_bg(slide, color=WHITE):
    bg = add_rect(slide, 0, 0, W, H, fill_color=color)
    return bg


def header_bar(slide, text, subtitle=None):
    add_rect(slide, 0, 0, W, Inches(1.2), fill_color=NAVY)
    add_text(slide, text, Inches(0.4), Inches(0.15), Inches(12), Inches(0.8),
             font_size=32, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, subtitle, Inches(0.4), Inches(0.85), Inches(12), Inches(0.4),
                 font_size=16, color=GOLD)


def footer(slide, text="Emport AI — 中小企業のためのAI活用入門"):
    add_rect(slide, 0, Inches(7.1), W, Inches(0.4), fill_color=NAVY)
    add_text(slide, text, Inches(0.3), Inches(7.12), Inches(10), Inches(0.3),
             font_size=12, color=WHITE)


def bullet_box(slide, lines, left, top, width, height, font_size=20, color=CHARCOAL, spacing=1.2):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.color.rgb = color


def card(slide, left, top, width, height, title, body_lines,
         bg=LIGHT_GRAY, title_color=NAVY, body_color=CHARCOAL):
    add_rect(slide, left, top, width, height, fill_color=bg)
    add_text(slide, title, left + Inches(0.15), top + Inches(0.1),
             width - Inches(0.3), Inches(0.5),
             font_size=18, bold=True, color=title_color)
    bullet_box(slide, body_lines,
               left + Inches(0.15), top + Inches(0.55),
               width - Inches(0.3), height - Inches(0.7),
               font_size=16, color=body_color)


# ─────────────────────────────────────────
# SLIDE 01: 表紙
# ─────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
slide_bg(s, NAVY)
add_rect(s, 0, Inches(2.8), W, Inches(2.2), fill_color=RGBColor(0x14, 0x28, 0x45))
add_text(s, "中小企業のためのAI活用入門", Inches(0.8), Inches(1.0), Inches(11.5), Inches(1.2),
         font_size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "— 明日から使える業務効率化のヒント —", Inches(0.8), Inches(2.1), Inches(11.5), Inches(0.6),
         font_size=24, color=GOLD, align=PP_ALIGN.CENTER)
add_rect(s, Inches(5.5), Inches(2.95), Inches(2.3), Inches(0.05), fill_color=GOLD)
add_text(s, "Emport AI", Inches(0.8), Inches(3.2), Inches(11.5), Inches(0.6),
         font_size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "主催：○○商工会議所", Inches(0.8), Inches(3.8), Inches(11.5), Inches(0.5),
         font_size=18, color=GOLD, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────
# SLIDE 02: 今日のゴール
# ─────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
slide_bg(s)
header_bar(s, "今日のゴール")
footer(s)
add_text(s, "今日終わったとき、こうなっていてほしい", Inches(0.5), Inches(1.4), Inches(12), Inches(0.5),
         font_size=22, bold=True, color=NAVY)
goals = [
    "✅  「AIって自分の会社に使えるかも」と思っている",
    "✅  「何から始めればいいか」がわかっている",
    "✅  「コストと補助金」のイメージがついている",
]
bullet_box(s, goals, Inches(0.8), Inches(2.0), Inches(11.5), Inches(3.0), font_size=24, color=CHARCOAL)
add_text(s, "※ AIの全部を知る必要はありません。今日は「最初の一歩」の話だけします。",
         Inches(0.8), Inches(5.8), Inches(11.5), Inches(0.5), font_size=16, color=RGBColor(0x88, 0x88, 0x88))

# ─────────────────────────────────────────
# SLIDE 03: AIの3つの誤解
# ─────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
slide_bg(s)
header_bar(s, "まず、よくある誤解を3つ解消します")
footer(s)
misconceptions = [
    ("① AIは大企業のもの", "→  × 誤解です"),
    ("② 導入に何百万もかかる", "→  × 誤解です"),
    ("③ ITが苦手な社員には無理", "→  × 誤解です"),
]
tops = [Inches(1.5), Inches(3.0), Inches(4.5)]
for (label, cross), top in zip(misconceptions, tops):
    add_rect(s, Inches(0.5), top, Inches(8.5), Inches(1.1), fill_color=LIGHT_GRAY)
    add_text(s, label, Inches(0.7), top + Inches(0.25), Inches(6.5), Inches(0.7),
             font_size=22, bold=True, color=CHARCOAL)
    add_text(s, cross, Inches(9.0), top + Inches(0.25), Inches(3.5), Inches(0.7),
             font_size=22, bold=True, color=RED)

add_text(s, "このあとひとつずつ見ていきます →", Inches(0.5), Inches(6.3), Inches(12), Inches(0.5),
         font_size=18, color=NAVY, align=PP_ALIGN.RIGHT)

# ─────────────────────────────────────────
# SLIDE 04: 誤解① 大企業のもの
# ─────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
slide_bg(s)
header_bar(s, "誤解①「AIは大企業のもの」", "→  これは間違いです")
footer(s)
add_rect(s, Inches(0.4), Inches(1.5), Inches(5.8), Inches(3.5), fill_color=RGBColor(0xFF, 0xEC, 0xEC))
add_text(s, "よくある誤解", Inches(0.6), Inches(1.6), Inches(5.5), Inches(0.5),
         font_size=16, bold=True, color=RED)
bullet_box(s, ["「トヨタや楽天がやるもの」", "「うちには関係ない」", "「専門家がいないと無理」"],
           Inches(0.6), Inches(2.1), Inches(5.5), Inches(2.5), font_size=20, color=CHARCOAL)

add_rect(s, Inches(6.8), Inches(1.5), Inches(6.0), Inches(3.5), fill_color=RGBColor(0xE8, 0xF5, 0xE9))
add_text(s, "事実", Inches(7.0), Inches(1.6), Inches(5.5), Inches(0.5),
         font_size=16, bold=True, color=GREEN)
add_text(s, "従業員20名以下の小規模事業者でも", Inches(7.0), Inches(2.1), Inches(5.5), Inches(0.5),
         font_size=18, color=CHARCOAL)
add_text(s, "3社に1社", Inches(7.0), Inches(2.65), Inches(5.5), Inches(0.9),
         font_size=44, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
add_text(s, "がすでにAIツールを活用", Inches(7.0), Inches(3.5), Inches(5.5), Inches(0.5),
         font_size=18, color=CHARCOAL, align=PP_ALIGN.CENTER)

add_rect(s, Inches(0.4), Inches(5.3), Inches(12.4), Inches(0.9), fill_color=NAVY)
add_text(s, "隣の競合が、もうこっそり使い始めているかもしれません。",
         Inches(0.6), Inches(5.4), Inches(12.0), Inches(0.7),
         font_size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────
# SLIDE 05: 誤解② コスト
# ─────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
slide_bg(s)
header_bar(s, "誤解②「導入に何百万もかかる」", "→  これも間違いです")
footer(s)
rows = [
    ("よくある誤解", "現実"),
    ("システム開発 ＝ 高い", "月額1,000円〜のツールがある"),
    ("自社専用が必要", "既製品で十分なことが多い"),
    ("補助金は面倒", "申請サポートがある"),
]
for i, (left_text, right_text) in enumerate(rows):
    top = Inches(1.4 + i * 1.1)
    bg = NAVY if i == 0 else (LIGHT_GRAY if i % 2 == 1 else WHITE)
    txt_color = WHITE if i == 0 else CHARCOAL
    add_rect(s, Inches(0.4), top, Inches(6.0), Inches(1.0), fill_color=bg)
    add_rect(s, Inches(6.8), top, Inches(6.0), Inches(1.0), fill_color=bg)
    add_text(s, left_text, Inches(0.6), top + Inches(0.2), Inches(5.8), Inches(0.7),
             font_size=18 if i > 0 else 16, bold=(i == 0), color=txt_color)
    add_text(s, right_text, Inches(7.0), top + Inches(0.2), Inches(5.8), Inches(0.7),
             font_size=18 if i > 0 else 16, bold=(i == 0), color=txt_color)
    add_text(s, "→", Inches(6.3), top + Inches(0.25), Inches(0.5), Inches(0.5),
             font_size=20, bold=True, color=GOLD if i == 0 else NAVY)

add_rect(s, Inches(0.4), Inches(6.0), Inches(12.4), Inches(0.9), fill_color=GOLD)
add_text(s, "IT導入補助金を活用すれば、導入費用の最大75%が補助される",
         Inches(0.6), Inches(6.1), Inches(12.0), Inches(0.7),
         font_size=20, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────
# SLIDE 06: 誤解③ ITが苦手
# ─────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
slide_bg(s)
header_bar(s, "誤解③「ITが苦手な社員には無理」", "→  今のAIは別物です")
footer(s)
add_text(s, "今のAIは「スマホのアプリ」と同じ感覚で使える", Inches(0.5), Inches(1.5), Inches(12), Inches(0.6),
         font_size=24, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
add_rect(s, Inches(1.5), Inches(2.3), Inches(10.0), Inches(1.8), fill_color=RGBColor(0xE8, 0xF5, 0xE9))
add_text(s, "実例：60代のパート従業員の方が2週間でAIツールを使いこなした",
         Inches(1.7), Inches(2.5), Inches(9.6), Inches(1.2),
         font_size=22, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
add_text(s, "難しいのはツールではなく\n「何に使うか」を決めること",
         Inches(1.5), Inches(4.4), Inches(10), Inches(1.4),
         font_size=26, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
add_text(s, "それを一緒に考えるのが私たちの仕事です。",
         Inches(1.5), Inches(5.9), Inches(10), Inches(0.6),
         font_size=20, color=CHARCOAL, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────
# SLIDE 07: 山口県の現実
# ─────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
slide_bg(s)
header_bar(s, "避けられない現実 — 山口県の人手不足")
footer(s)
facts = [
    "山口県の有効求人倍率：全国平均を上回り続けている",
    "2030年までに生産年齢人口がさらに10%以上減少する見込み",
    "求人を出しても来ない。来ても続かない。",
]
bullet_box(s, facts, Inches(0.6), Inches(1.6), Inches(11.5), Inches(2.5), font_size=22, color=CHARCOAL)
add_rect(s, Inches(0.5), Inches(4.3), Inches(12.2), Inches(1.2), fill_color=NAVY)
add_text(s, "5年後、今のやり方で会社を回せますか？",
         Inches(0.7), Inches(4.5), Inches(11.8), Inches(0.8),
         font_size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "AIは「人を減らす」道具ではなく\n「今いる人でもっとできる」道具です。",
         Inches(0.5), Inches(5.7), Inches(12.2), Inches(1.2),
         font_size=22, color=CHARCOAL, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────
# SLIDE 08: 事例 製造業
# ─────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
slide_bg(s)
header_bar(s, "業種別 活用事例 — 製造業")
footer(s)
card(s, Inches(0.4), Inches(1.5), Inches(6.0), Inches(2.5),
     "【部品メーカー / 従業員15名】",
     ["課題: 検品がベテラン2名に依存",
      "解決: AIカメラによる外観検査を導入",
      "結果: ✅ 品質水準を維持しつつ体制を強化",
      "費用: 150万円 → 補助金100万円 → 実質50万円"])
card(s, Inches(6.8), Inches(1.5), Inches(6.0), Inches(2.5),
     "【食品製造 / 従業員30名】",
     ["課題: FAX受発注処理に毎日2時間",
      "解決: AIによるFAX自動読み取り・データ化",
      "結果: ✅ 処理時間 2時間 → 20分（1/6）",
      "費用: 月額2万円（補助金で初期費用カバー）"])
add_rect(s, Inches(0.4), Inches(4.3), Inches(12.4), Inches(0.8), fill_color=GOLD)
add_text(s, "書類・データ処理が発生する業種なら、ほぼすべてで活用できます",
         Inches(0.6), Inches(4.4), Inches(12.0), Inches(0.6),
         font_size=20, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────
# SLIDE 09: 事例 サービス業
# ─────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
slide_bg(s)
header_bar(s, "業種別 活用事例 — サービス業・小売")
footer(s)
card(s, Inches(0.4), Inches(1.5), Inches(6.0), Inches(2.8),
     "【旅館 / 従業員20名】",
     ["課題: 予約問い合わせ対応に毎日3〜4時間",
      "解決: AIチャットボットをHP・LINEに設置",
      "結果: ✅ 問い合わせの75%をAIが自動対応",
      "       ✅ 深夜予約も自動化。機会損失ゼロ"])
card(s, Inches(6.8), Inches(1.5), Inches(6.0), Inches(2.8),
     "【小売店 / 従業員8名】",
     ["課題: 発注作業に時間がかかり在庫の過不足頻発",
      "解決: 過去データをAIが分析し発注量を自動提案",
      "結果: ✅ 在庫ロス30%削減",
      "       ✅ 発注作業時間が半分に短縮"])

# ─────────────────────────────────────────
# SLIDE 10: ワーク
# ─────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
slide_bg(s, LIGHT_GRAY)
header_bar(s, "3分ワーク：自社に当てはめてみよう ✏️")
footer(s)
add_text(s, "お手元のシートに書いてください", Inches(0.5), Inches(1.4), Inches(12), Inches(0.5),
         font_size=20, color=CHARCOAL)
questions = [
    "Q1.  今の業務で「時間がかかって困っている」作業は何ですか？",
    "Q2.  「誰かに任せたい」と思っている作業はありますか？",
    "Q3.  作業時間が半分になったら、何に使いたいですか？",
]
for i, q in enumerate(questions):
    top = Inches(2.1 + i * 1.4)
    add_rect(s, Inches(0.5), top, Inches(12.2), Inches(1.1), fill_color=WHITE)
    add_text(s, q, Inches(0.7), top + Inches(0.2), Inches(11.8), Inches(0.8),
             font_size=20, color=NAVY, bold=True)
add_text(s, "※ 正解はありません。思いつくままで大丈夫です。",
         Inches(0.5), Inches(6.7), Inches(12), Inches(0.4),
         font_size=14, color=RGBColor(0x88, 0x88, 0x88))

# ─────────────────────────────────────────
# SLIDE 11: 補助金
# ─────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
slide_bg(s)
header_bar(s, "使える補助金 — IT導入補助金（経済産業省）")
footer(s)
specs = [
    ("対象", "中小企業・小規模事業者"),
    ("補助率", "最大 75%"),
    ("補助額", "5万円 〜 450万円"),
    ("対象経費", "ITツール・クラウド・導入支援費"),
]
for i, (label, val) in enumerate(specs):
    top = Inches(1.5 + i * 1.1)
    add_rect(s, Inches(0.4), top, Inches(3.0), Inches(0.9), fill_color=NAVY)
    add_rect(s, Inches(3.5), top, Inches(9.3), Inches(0.9), fill_color=LIGHT_GRAY)
    add_text(s, label, Inches(0.5), top + Inches(0.15), Inches(2.8), Inches(0.65),
             font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, val, Inches(3.7), top + Inches(0.15), Inches(9.0), Inches(0.65),
             font_size=20, color=CHARCOAL)

add_rect(s, Inches(0.4), Inches(6.0), Inches(12.4), Inches(0.9), fill_color=GOLD)
add_text(s, "100万円の導入 → 補助75万円 → 自己負担 25万円",
         Inches(0.6), Inches(6.1), Inches(12.0), Inches(0.7),
         font_size=22, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────
# SLIDE 12: AI業務診断
# ─────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
slide_bg(s)
header_bar(s, "最初の一歩：AI業務診断")
footer(s)
add_rect(s, Inches(0.4), Inches(1.4), Inches(12.4), Inches(1.0), fill_color=NAVY)
add_text(s, "私たちが勧める「最初の一歩」", Inches(0.6), Inches(1.5), Inches(12.0), Inches(0.7),
         font_size=26, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
steps = [
    ("半日", "御社に訪問してヒアリング"),
    ("1週間後", "AIレポートをお渡し"),
    ("報告会", "結果をもとに次の方針を一緒に決める"),
]
for i, (timing, desc) in enumerate(steps):
    left = Inches(0.5 + i * 4.3)
    add_rect(s, left, Inches(2.7), Inches(3.9), Inches(1.8), fill_color=LIGHT_GRAY)
    add_text(s, timing, left + Inches(0.1), Inches(2.8), Inches(3.7), Inches(0.6),
             font_size=22, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    add_text(s, desc, left + Inches(0.1), Inches(3.35), Inches(3.7), Inches(0.9),
             font_size=17, color=CHARCOAL, align=PP_ALIGN.CENTER)

report_items = [
    "✅  AIで効率化できる業務リスト",
    "✅  おすすめツール・ソリューション案",
    "✅  費用と補助金の活用シミュレーション",
    "✅  最初に手をつけるべき優先順位",
]
bullet_box(s, report_items, Inches(0.6), Inches(4.8), Inches(11.5), Inches(2.0),
           font_size=20, color=CHARCOAL)

add_rect(s, Inches(0.4), Inches(6.3), Inches(12.4), Inches(0.7), fill_color=GOLD)
add_text(s, "費用: 5万円〜10万円　|　診断だけで終わってもOK",
         Inches(0.6), Inches(6.4), Inches(12.0), Inches(0.5),
         font_size=20, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────
# SLIDE 13: まとめ
# ─────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
slide_bg(s)
header_bar(s, "本日のまとめ")
footer(s)
summaries = [
    "① AIは中小企業にこそ使える（3つの誤解を解消）",
    "② 山口県の人手不足に、AIは現実的な解決策",
    "③ 製造・旅館・小売・建設で具体的な実績あり",
    "④ 補助金を使えば自己負担を最小化できる",
    "⑤ 最初の一歩は「AI業務診断」（5〜10万円）",
]
bullet_box(s, summaries, Inches(0.6), Inches(1.5), Inches(12.0), Inches(4.0),
           font_size=22, color=CHARCOAL)
add_rect(s, Inches(0.4), Inches(5.8), Inches(12.4), Inches(1.2), fill_color=NAVY)
add_text(s, "AIの波はもう来ている。問題はいつ乗るか。",
         Inches(0.6), Inches(6.0), Inches(12.0), Inches(0.8),
         font_size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────
# SLIDE 14: 次のステップ
# ─────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
slide_bg(s)
header_bar(s, "今日ここから動ける3つの選択肢")
footer(s)
options = [
    ("① まず自分で試す", "ChatGPT無料版を1週間使ってみる\nメール下書き・議事録・企画書の叩き台から"),
    ("② 相談だけする（無料）", "セミナー終了後にスタッフに声をかける\n名刺をお渡しします"),
    ("③ AI業務診断を申し込む", "今日お申し込みの方は優先対応します\nまずはお気軽にどうぞ"),
]
for i, (title, desc) in enumerate(options):
    top = Inches(1.5 + i * 1.7)
    bg = NAVY if i == 2 else LIGHT_GRAY
    tc = WHITE if i == 2 else NAVY
    dc = WHITE if i == 2 else CHARCOAL
    add_rect(s, Inches(0.5), top, Inches(12.2), Inches(1.4), fill_color=bg)
    add_text(s, title, Inches(0.7), top + Inches(0.1), Inches(12.0), Inches(0.5),
             font_size=22, bold=True, color=tc)
    add_text(s, desc, Inches(0.7), top + Inches(0.6), Inches(12.0), Inches(0.7),
             font_size=17, color=dc)

# ─────────────────────────────────────────
# SLIDE 15: 終了スライド
# ─────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
slide_bg(s, NAVY)
add_rect(s, 0, Inches(2.8), W, Inches(2.0), fill_color=RGBColor(0x14, 0x28, 0x45))
add_text(s, "ご参加ありがとうございました", Inches(0.5), Inches(1.2), Inches(12.3), Inches(1.2),
         font_size=38, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "「今日が、御社のAI活用のスタートラインになれば嬉しいです。」",
         Inches(0.5), Inches(3.0), Inches(12.3), Inches(0.9),
         font_size=22, color=GOLD, align=PP_ALIGN.CENTER)
add_text(s, "— Emport AI  CEO アレン", Inches(0.5), Inches(4.0), Inches(12.3), Inches(0.6),
         font_size=20, color=WHITE, align=PP_ALIGN.CENTER)
add_rect(s, Inches(5.0), Inches(4.9), Inches(3.3), Inches(0.06), fill_color=GOLD)
add_text(s, "終了後、個別相談受付中。お気軽に声をかけてください。",
         Inches(0.5), Inches(5.2), Inches(12.3), Inches(0.6),
         font_size=18, color=RGBColor(0xCC, 0xCC, 0xCC), align=PP_ALIGN.CENTER)

# 保存
out_path = "departments/sales/emport-ai-seminar-slides.pptx"
prs.save(out_path)
print("Slides generated: " + out_path)
print("Total slides: " + str(len(prs.slides)))
